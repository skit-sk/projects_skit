#!/usr/bin/env python3
"""
Генерация 30-слайдовой презентации учебника институционального трейдинга.
Данные: projects/01_fundament_rf/data/card/
Выход: docs/trading/Презентация_учебника.pptx + .pdf
"""

import json, os, math, io, warnings
from pathlib import Path
from datetime import datetime, timedelta
from collections import defaultdict

import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from matplotlib.patches import Rectangle, FancyBboxPatch, FancyArrowPatch
import matplotlib.patches as mpatches
from matplotlib.path import Path as MPath

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')

# ─── PATHS ────────────────────────────────────────────
PROJECT_ROOT = Path('/home/user_aioc/workspace')
CARD_DIR = PROJECT_ROOT / 'projects/01_fundament_rf/data/card'
OUTPUT_DIR = PROJECT_ROOT / 'docs/trading'
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ─── DARK THEME ───────────────────────────────────────
BG_COLOR      = '#1a1a2e'
BG2_COLOR     = '#16213e'
TEXT_COLOR    = '#ffffff'
TEXT_DIM      = '#cccccc'
BULL_COLOR    = '#00c853'
BEAR_COLOR    = '#ff1744'
ACCENT1       = '#ff9100'  # POC / OB
ACCENT2       = '#00e5ff'  # FVG
ACCENT3       = '#ffea00'  # сигналы
ACCENT4       = '#7c4dff'  # фиолетовый
VOLUME_COLOR  = '#6495ed'
OB_BULL       = '#00c853'
OB_BEAR       = '#ff1744'
FVG_FILL      = '#00e5ff'
OTE_FILL      = '#00e5ff'

# mpl rcParams
plt.rcParams.update({
    'figure.facecolor': BG_COLOR,
    'axes.facecolor': BG2_COLOR,
    'axes.edgecolor': '#444466',
    'axes.labelcolor': '#ffffff',
    'axes.titlecolor': '#ffffff',
    'xtick.color': '#999999',
    'ytick.color': '#999999',
    'grid.color': '#333355',
    'grid.alpha': 0.3,
    'text.color': '#ffffff',
    'font.family': 'DejaVu Sans',
    'font.size': 10,
})

# ─── DATA LOADERS ─────────────────────────────────────
def load_1d(symbol):
    """Загрузить обогащённые 1D данные по символу."""
    for p in CARD_DIR.glob(f'{symbol}_*/**/*_1D.json'):
        with open(p) as f:
            return json.load(f)
    return None

def load_raw(symbol):
    for p in CARD_DIR.glob(f'{symbol}_*/**/*_RAW.json'):
        with open(p) as f:
            return json.load(f)
    return None

def load_card(symbol):
    for p in CARD_DIR.glob(f'{symbol}_*/*.json'):
        if '_1D' not in p.name and '_RAW' not in p.name and '_1h' not in p.name and '_4h' not in p.name:
            with open(p) as f:
                return json.load(f)
    return None

def get_candles(data):
    return data.get('candles', [])

def candle_array(candles, field):
    return np.array([c[field] for c in candles], dtype=float)

def dates_array(candles):
    return [datetime.strptime(c['date'], '%Y-%m-%d') for c in candles]

def price_format(val):
    if val < 0.01: return f'{val:.6f}'
    if val < 1: return f'{val:.4f}'
    if val < 100: return f'{val:.2f}'
    return f'{val:.1f}'

# ─── DARK STYLING HELPERS ─────────────────────────────
def style_ax(ax, title='', xlabel='', ylabel=''):
    ax.set_facecolor(BG2_COLOR)
    ax.set_title(title, color=TEXT_COLOR, fontsize=13, fontweight='bold', pad=12)
    ax.set_xlabel(xlabel, color=TEXT_DIM)
    ax.set_ylabel(ylabel, color=TEXT_DIM)
    ax.tick_params(colors=TEXT_DIM)
    ax.grid(True, alpha=0.15, color='#444466')
    for spine in ax.spines.values():
        spine.set_color('#444466')
    return ax

def add_info_box(ax, text, x=0.02, y=0.96, fontsize=9, color=TEXT_DIM):
    ax.text(x, y, text, transform=ax.transAxes, fontsize=fontsize,
            color=color, va='top', ha='left',
            bbox=dict(facecolor='#111122', alpha=0.7, edgecolor='none', pad=4))

def draw_candlestick(ax, candles, start=0, end=None, width=0.6):
    """Рисует свечной график."""
    if end is None: end = len(candles)
    subset = candles[start:end]
    dates = dates_array(subset)
    opens = candle_array(subset, 'open')
    highs = candle_array(subset, 'high')
    lows = candle_array(subset, 'low')
    closes = candle_array(subset, 'close')
    dnums = matplotlib.dates.date2num(dates)

    for i in range(len(subset)):
        color = BULL_COLOR if closes[i] >= opens[i] else BEAR_COLOR
        ax.plot([dnums[i], dnums[i]], [lows[i], highs[i]], color=color, linewidth=0.8, alpha=0.7)
        w = width * 0.4
        rect = Rectangle((dnums[i]-w, min(opens[i], closes[i])), 2*w, abs(closes[i]-opens[i]),
                         facecolor=color, edgecolor=color, linewidth=0.5, alpha=0.9)
        ax.add_patch(rect)
    return dnums

# ─── INDICATOR CALCULATIONS ───────────────────────────
def calc_sma(data, period):
    if len(data) < period: return np.full_like(data, np.nan)
    res = np.full_like(data, np.nan)
    for i in range(period-1, len(data)):
        res[i] = np.mean(data[i-period+1:i+1])
    return res

def calc_ema(data, period):
    if len(data) < period: return np.full_like(data, np.nan)
    res = np.full_like(data, np.nan)
    multiplier = 2 / (period + 1)
    res[period-1] = np.mean(data[:period])
    for i in range(period, len(data)):
        res[i] = (data[i] - res[i-1]) * multiplier + res[i-1]
    return res

def calc_atr(candles, period=14):
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    closes = candle_array(candles, 'close')
    tr = np.zeros(len(candles))
    tr[0] = highs[0] - lows[0]
    for i in range(1, len(candles)):
        tr[i] = max(highs[i]-lows[i], abs(highs[i]-closes[i-1]), abs(lows[i]-closes[i-1]))
    return calc_ema(tr, period)

def calc_adx(candles, period=14):
    high = candle_array(candles, 'high')
    low = candle_array(candles, 'low')
    close = candle_array(candles, 'close')
    up = np.zeros(len(candles))
    down = np.zeros(len(candles))
    for i in range(1, len(candles)):
        up[i] = high[i] - high[i-1]
        down[i] = low[i-1] - low[i]
    atr = calc_atr(candles, period)
    di_plus = np.full(len(candles), np.nan)
    di_minus = np.full(len(candles), np.nan)
    for i in range(1, len(candles)):
        if atr[i] != 0 and not np.isnan(atr[i]):
            di_plus[i] = 100 * (up[i] if up[i] > down[i] and up[i] > 0 else 0) / atr[i]
            di_minus[i] = 100 * (down[i] if down[i] > up[i] and down[i] > 0 else 0) / atr[i]
    dx = np.full(len(candles), np.nan)
    for i in range(1, len(candles)):
        s = di_plus[i] + di_minus[i]
        if s != 0 and not np.isnan(s):
            dx[i] = 100 * abs(di_plus[i] - di_minus[i]) / s
    return calc_ema(dx, period)

def calc_macd(candles):
    close = candle_array(candles, 'close')
    ema12 = calc_ema(close, 12)
    ema26 = calc_ema(close, 26)
    macd_line = ema12 - ema26
    signal = calc_ema(macd_line, 9)
    hist = macd_line - signal
    return macd_line, signal, hist

def calc_rsi(candles, period=14):
    close = candle_array(candles, 'close')
    delta = np.diff(close)
    gain = np.where(delta > 0, delta, 0)
    loss = np.where(delta < 0, -delta, 0)
    avg_gain = np.full(len(close), np.nan)
    avg_loss = np.full(len(close), np.nan)
    avg_gain[period] = np.mean(gain[:period])
    avg_loss[period] = np.mean(loss[:period])
    for i in range(period+1, len(close)):
        avg_gain[i] = (avg_gain[i-1]*13 + gain[i-1]) / period
        avg_loss[i] = (avg_loss[i-1]*13 + loss[i-1]) / period
    rs = avg_gain / np.where(avg_loss==0, 0.001, avg_loss)
    rsi = 100 - 100 / (1 + rs)
    return rsi

def calc_bb(candles, period=20, std=2):
    close = candle_array(candles, 'close')
    sma = calc_sma(close, period)
    upper = np.full(len(close), np.nan)
    lower = np.full(len(close), np.nan)
    for i in range(period-1, len(close)):
        s = np.std(close[i-period+1:i+1])
        upper[i] = sma[i] + std * s
        lower[i] = sma[i] - std * s
    return sma, upper, lower

# ─── SLIDE GENERATORS ─────────────────────────────────
def make_slide1_adx():
    """Слайд 1: ADX-фильтр <30 vs >30"""
    data = load_1d('ATOM')
    if not data: return None
    candles = get_candles(data)[-200:]
    if len(candles) < 30: return None

    closes = candle_array(candles, 'close')
    dates = dates_array(candles)
    adx = calc_adx(candles, 14)
    adx[np.isnan(adx)] = 0  # replace NaN

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(16, 8), gridspec_kw={'height_ratios': [3, 1]})
    style_ax(ax1, 'ADX-фильтр: определение режима рынка (ATOM)', '')
    style_ax(ax2, 'ADX(14) — разделение на тренд/боковик', '', 'ADX')

    dnums = draw_candlestick(ax1, candles)
    dnums_arr = np.array([mdates.date2num(d) for d in dates_array(candles)])

    # Color zones
    for i in range(len(candles)):
        if i < 14: continue
        if adx[i] > 30:
            ax1.axvspan(dnums_arr[i]-0.5, dnums_arr[i]+0.5, alpha=0.08, color=ACCENT1)
        else:
            ax1.axvspan(dnums_arr[i]-0.5, dnums_arr[i]+0.5, alpha=0.05, color=ACCENT2)

    # ADX line
    ax2.plot(dnums, adx, color='#ff9100', linewidth=1.8, label='ADX(14)')
    ax2.axhline(y=30, color=BEAR_COLOR, linestyle='--', linewidth=1, alpha=0.7, label='ADX=30 (порог)')
    ax2.axhline(y=25, color=TEXT_DIM, linestyle=':', linewidth=0.8, alpha=0.5)

    ax2.fill_between(dnums, 0, adx, where=(adx>30), color=ACCENT1, alpha=0.15)
    ax2.fill_between(dnums, 0, np.where(adx<30, 30, np.nan), color=ACCENT2, alpha=0.1)
    ax2.legend(loc='upper right', fontsize=9, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)
    adx_max = np.nanmax(adx[14:])
    if np.isnan(adx_max) or adx_max < 40: adx_max = 40
    ax2.set_ylim(0, adx_max * 1.1)

    add_info_box(ax1, 'ADX < 30 → Mean Reversion (Bollinger, Stochastic)  |  ADX > 30 → Тренд (MACD, SMC)')
    fig.tight_layout()
    return fig

def make_slide2_macd_stoch():
    """Слайд 2: MACD + Stochastic"""
    data = load_1d('ETH')
    if not data: return None
    candles = get_candles(data)
    if len(candles) < 26: return None

    closes = candle_array(candles, 'close')
    dates = dates_array(candles)
    macd_line, signal, hist = calc_macd(candles)
    dnums = dates_array(candles)

    # Stochastic
    high14 = candle_array(candles, 'high')
    low14 = candle_array(candles, 'low')
    stoch_k = np.full(len(candles), np.nan)
    for i in range(13, len(candles)):
        h = np.max(high14[i-13:i+1])
        l = np.min(low14[i-13:i+1])
        if h - l != 0:
            stoch_k[i] = (closes[i] - l) / (h - l) * 100

    fig, (ax1, ax2, ax3) = plt.subplots(3, 1, figsize=(16, 9), gridspec_kw={'height_ratios': [3, 1.2, 1.2]})
    style_ax(ax1, 'MACD + Stochastic — синхронизация индикаторов (ETH)', '')
    style_ax(ax2, 'MACD (12, 26, 9)', '', '')
    style_ax(ax3, 'Stochastic %K(14, 3)', '', '')

    draw_candlestick(ax1, candles)

    # MACD
    ax2.plot(dnums, macd_line, color=ACCENT2, linewidth=1.5, label='MACD Line')
    ax2.plot(dnums, signal, color=ACCENT3, linewidth=1.5, label='Signal Line')
    for i in range(len(candles)):
        if not np.isnan(hist[i]):
            color = BULL_COLOR if hist[i] >= 0 else BEAR_COLOR
            ax2.bar(dnums[i], hist[i], width=0.6, color=color, alpha=0.6)
    ax2.axhline(y=0, color=TEXT_DIM, linewidth=0.5)
    ax2.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # Stochastic
    ax3.plot(dnums, stoch_k, color='#7c4dff', linewidth=1.5, label='%K')
    ax3.axhline(y=80, color=BEAR_COLOR, linestyle='--', alpha=0.5)
    ax3.axhline(y=20, color=BULL_COLOR, linestyle='--', alpha=0.5)
    ax3.fill_between(dnums, 80, 100, alpha=0.08, color=BEAR_COLOR)
    ax3.fill_between(dnums, 0, 20, alpha=0.08, color=BULL_COLOR)
    ax3.set_ylim(-5, 105)
    ax3.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    add_info_box(ax1, 'MACD Line = EMA₁₂ − EMA₂₆  |  Signal = EMA₉  |  Stochastic %K = (C−L₁₄)/(H₁₄−L₁₄)×100')
    fig.tight_layout()
    return fig

def make_slide3_bb_atr():
    """Слайд 3: Bollinger + ATR"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)[-100:]
    if len(candles) < 20: return None

    closes = candle_array(candles, 'close')
    dates = dates_array(candles)
    sma, upper, lower = calc_bb(candles, 20, 2)
    atr = calc_atr(candles, 14)
    dnums = dates_array(candles)

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(16, 8), gridspec_kw={'height_ratios': [3, 1]})
    style_ax(ax1, 'Bollinger Bands (20, 2) + ATR(14) — волатильность (API3)', '')
    style_ax(ax2, 'ATR(14) — Average True Range', '', '')

    draw_candlestick(ax1, candles)

    # Bollinger bands
    ax1.plot(dnums, sma, color=ACCENT3, linewidth=1, alpha=0.8, label='SMA(20)')
    ax1.plot(dnums, upper, color=ACCENT2, linewidth=0.8, alpha=0.6, label='Upper BB')
    ax1.plot(dnums, lower, color=ACCENT2, linewidth=0.8, alpha=0.6, label='Lower BB')
    ax1.fill_between(dnums, upper, lower, alpha=0.06, color=ACCENT2)
    ax1.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # ATR
    ax2.fill_between(dnums, 0, atr, color='#7c4dff', alpha=0.4)
    ax2.plot(dnums, atr, color='#7c4dff', linewidth=1.8)
    ax2.axhline(y=np.nanmean(atr[13:]), color=ACCENT3, linestyle='--', alpha=0.5, label=f'средн.')
    ax2.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    add_info_box(ax1, 'Bollinger: SMA₂₀ ± 2σ  |  ATR: True Range → EMA₁₄  |  DynamicStep = ATR×0.2')
    fig.tight_layout()
    return fig

def make_slide4_candle_patterns():
    """Слайд 4: Свечные паттерны (Doji, Hammer, Engulfing)"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)[-80:]
    if len(candles) < 10: return None

    closes = candle_array(candles, 'close')
    opens = candle_array(candles, 'open')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    dnums = dates_array(candles)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))
    style_ax(ax1, 'Свечные паттерны на графике (API3)', '')
    style_ax(ax2, 'Примеры ключевых паттернов', '')

    # Main chart
    draw_candlestick(ax1, candles)

    # Find and mark patterns
    annotations = []
    for i in range(1, len(candles)-1):
        body = abs(closes[i] - opens[i])
        total_range = highs[i] - lows[i]
        if total_range == 0: continue
        upper_wick = highs[i] - max(opens[i], closes[i])
        lower_wick = min(opens[i], closes[i]) - lows[i]
        body_pct = body / total_range

        # Doji: body < 5% of range
        if body_pct < 0.05 and total_range > 0:
            ax1.annotate('Doji', (dnums[i], highs[i]), textcoords='offset points',
                        xytext=(0, 8), fontsize=7, color=ACCENT3, ha='center',
                        arrowprops=dict(arrowstyle='->', color=ACCENT3, lw=0.5))
        # Hammer: lower wick > 2*body, body in upper half
        if lower_wick > 2 * body and body_pct < 0.3 and closes[i] >= opens[i]:
            ax1.annotate('Hammer', (dnums[i], lows[i]), textcoords='offset points',
                        xytext=(0, -10), fontsize=7, color=BULL_COLOR, ha='center',
                        arrowprops=dict(arrowstyle='->', color=BULL_COLOR, lw=0.5))
        # Engulfing
        if i > 0:
            prev_body = abs(closes[i-1] - opens[i-1])
            if closes[i] > opens[i] and closes[i-1] < opens[i-1]:
                if closes[i] > opens[i-1] and opens[i] < closes[i-1]:
                    ax1.annotate('Bull Engulf', (dnums[i], highs[i]), textcoords='offset points',
                                xytext=(0, 12), fontsize=6, color=BULL_COLOR, ha='center',
                                arrowprops=dict(arrowstyle='->', color=BULL_COLOR, lw=0.5))
            if closes[i] < opens[i] and closes[i-1] > opens[i-1]:
                if closes[i] < opens[i-1] and opens[i] > closes[i-1]:
                    ax1.annotate('Bear Engulf', (dnums[i], highs[i]), textcoords='offset points',
                                xytext=(0, 12), fontsize=6, color=BEAR_COLOR, ha='center',
                                arrowprops=dict(arrowstyle='->', color=BEAR_COLOR, lw=0.5))

    # Right panel: pattern examples
    patterns = [
        ('Doji', 'Open ≈ Close\nНеопределённость\nРазворот'),
        ('Hammer', 'Маленькое тело↑\nДлинная нижн. тень\nБычий разворот'),
        ('Shooting Star', 'Маленькое тело↓\nДлинная верхн. тень\nМедвежий разворот'),
        ('Bull Engulfing', 'Бычье тело\nпоглощает\nпредыдущее'),
        ('Morning Star', 'Медвежья→Doji→\n→Бычья свеча\nРазворот ↑'),
    ]
    y_pos = 0.9
    for name, desc in patterns:
        ax2.text(0.1, y_pos, f'◆ {name}', fontsize=11, fontweight='bold', color=ACCENT3, transform=ax2.transAxes)
        ax2.text(0.35, y_pos-0.02, desc, fontsize=9, color=TEXT_DIM, transform=ax2.transAxes)
        y_pos -= 0.18
    ax2.text(0.1, 0.02, 'Правило: валидация только на\nзначимом уровне (OB/S/R)', fontsize=9,
             color=TEXT_DIM, transform=ax2.transAxes, style='italic')
    ax2.axis('off')

    fig.tight_layout()
    return fig

def make_slide5_triple_star():
    """Слайд 5: Тройные паттерны (Morning/Evening Star)"""
    data = load_1d('ATOM')
    if not data: return None
    candles = get_candles(data)[-100:]
    if len(candles) < 5: return None

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))
    style_ax(ax1, 'Тройные разворотные паттерны (ATOM)', '')
    draw_candlestick(ax1, candles)

    closes = candle_array(candles, 'close')
    opens = candle_array(candles, 'open')
    highs = candle_array(candles, 'high')
    dnums = dates_array(candles)

    # Find potential 3-candle patterns
    for i in range(2, len(candles)):
        b1, b2, b3 = abs(closes[i-2]-opens[i-2]), abs(closes[i-1]-opens[i-1]), abs(closes[i]-opens[i])
        # Evening Star: big green, small body, big red
        if (closes[i-2] > opens[i-2] and abs(closes[i-2]-opens[i-2]) > 0.05 and
            b2 < b1 * 0.4 and closes[i] < opens[i] and closes[i] < opens[i-2]):
            ax1.annotate('Evening\nStar', (dnums[i], highs[i]), fontsize=7, color=BEAR_COLOR, ha='center',
                        xytext=(0, 10), textcoords='offset points',
                        arrowprops=dict(arrowstyle='->', color=BEAR_COLOR, lw=0.5))
        # Morning Star: big red, small body, big green
        if (closes[i-2] < opens[i-2] and abs(closes[i-2]-opens[i-2]) > 0.05 and
            b2 < b1 * 0.4 and closes[i] > opens[i] and closes[i] > opens[i-2]):
            ax1.annotate('Morning\nStar', (dnums[i], highs[i]), fontsize=7, color=BULL_COLOR, ha='center',
                        xytext=(0, 10), textcoords='offset points',
                        arrowprops=dict(arrowstyle='->', color=BULL_COLOR, lw=0.5))

    # Pattern diagrams
    ax2.axis('off')
    patterns_data = [
        ('Morning Star (↑)', ['Медвежья', 'Doji/Spinning', 'Бычья >50% 1-й'], BULL_COLOR),
        ('Evening Star (↓)', ['Бычья', 'Doji/Spinning', 'Медвежья <50% 1-й'], BEAR_COLOR),
        ('3 White Soldiers', ['3 бычьих подряд', 'Каждая выше хая', 'предыдущей'], BULL_COLOR),
        ('3 Black Crows', ['3 медвежьих подряд', 'Каждая ниже лоя', 'предыдущей'], BEAR_COLOR),
        ('3 Inside Up', ['Медвежья→Бычья', 'внутри→Третья', 'выше'], BULL_COLOR),
    ]
    y = 0.92
    for name, lines, clr in patterns_data:
        ax2.text(0.05, y, f'◆ {name}', fontsize=11, fontweight='bold', color=clr, transform=ax2.transAxes)
        for j, line in enumerate(lines):
            ax2.text(0.35, y-0.03-0.04*j, line, fontsize=8, color=TEXT_DIM, transform=ax2.transAxes)
        y -= 0.16

    fig.tight_layout()
    return fig

def make_slide6_graphical_patterns():
    """Слайд 6: H&S, Double Top, Flags"""
    data = load_1d('ATOM')
    if not data: return None
    candles = get_candles(data)[-200:]
    if len(candles) < 30: return None

    closes = candle_array(candles, 'close')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    dnums = dates_array(candles)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))
    style_ax(ax1, 'Графические паттерны: Голова и Плечи (ATOM)', '')
    draw_candlestick(ax1, candles)

    # Draw neckline for H&S (schematic)
    d_range = 50
    if len(candles) > d_range:
        # Left shoulder area
        mid = len(candles) - d_range
        ax1.axhline(y=np.mean(highs[mid:mid+10]), xmin=0.1, xmax=0.9,
                    color=ACCENT3, linestyle='--', linewidth=1, alpha=0.5)

    ax2.axis('off')
    patterns = [
        ('Разворотные', [
            ('Head & Shoulders', 'Голова между двух\nплеч, пробой шеи'),
            ('Double Top/Bottom', 'Два уровня = разворот\nцель = высота'),
            ('Triple Top/Bottom', 'Три касания уровня'),
            ('Rounding Top/Bottom', 'Плавная вершина/дно'),
        ]),
        ('Продолжения', [
            ('Flags/Pennants', 'Консолидация после\nимпульса →breakout'),
            ('Wedges', 'Клин: Rising(Bear)\nFalling(Bull)'),
            ('Cup & Handle', 'Округлое дно + ручка\n→бычий пробой'),
            ('Measured Move', 'Движение=откат\n=равное движение'),
        ]),
    ]
    y = 0.95
    for cat, items in patterns:
        ax2.text(0.05, y, f'══ {cat} ══', fontsize=11, fontweight='bold', color=ACCENT2, transform=ax2.transAxes)
        y -= 0.04
        for name, desc in items:
            ax2.text(0.08, y, f'▸ {name}', fontsize=9, color=ACCENT3, transform=ax2.transAxes)
            ax2.text(0.35, y-0.02, desc, fontsize=8, color=TEXT_DIM, transform=ax2.transAxes)
            y -= 0.09
        y -= 0.04

    add_info_box(ax1, 'Измерение паттерна: цель = высота конструкции, отложенная от точки пробоя')
    fig.tight_layout()
    return fig

def make_slide7_wyckoff():
    """Слайд 7: Wyckoff — 5 фаз A–E"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)
    if len(candles) < 50: return None

    closes = candle_array(candles, 'close')
    volumes = candle_array(candles, 'volume')
    dnums_raw = dates_array(candles)
    dnums_num = matplotlib.dates.date2num(dnums_raw)

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(16, 8), gridspec_kw={'height_ratios': [3, 1]})
    style_ax(ax1, 'Wyckoff: Цикл накопления/распределения (API3)', '')
    style_ax(ax2, 'Объём — подтверждение фаз', '', 'Volume')

    draw_candlestick(ax1, candles)

    # Mark phases visually
    n = len(candles)
    phases = [
        (0, int(n*0.15), 'A', '#ff9100'),
        (int(n*0.15), int(n*0.4), 'B', '#7c4dff'),
        (int(n*0.4), int(n*0.55), 'C (Spring)', '#00c853'),
        (int(n*0.55), int(n*0.75), 'D (SOS)', '#00e5ff'),
        (int(n*0.75), n-1, 'E (Markup)', '#ff1744'),
    ]
    for s, e, label, clr in phases:
        if e >= len(dnums_num): e = len(dnums_num)-1
        mid = (s + e) // 2
        if mid < len(dnums_num):
            ax1.axvspan(dnums_num[s]-0.5 if s>0 else dnums_num[0]-0.5, dnums_num[e]+0.5, alpha=0.06, color=clr)
            ax1.text(dnums_num[mid], ax1.get_ylim()[1], f'  {label}  ', fontsize=8, fontweight='bold',
                    color=clr, ha='center', va='bottom',
                    bbox=dict(facecolor='#111122', alpha=0.7, edgecolor=clr, pad=2))

    # Volume bars
    ax2.bar(dnums_raw, volumes, color=VOLUME_COLOR, width=0.6)
    ax2.set_yscale('log')

    add_info_box(ax1, 'Закон причины и следствия: боковик (причина) → импульс (следствие)')
    fig.tight_layout()
    return fig

def make_slide8_wyckoff_smc():
    """Слайд 8: Wyckoff ↔ SMC equivalence"""
    fig, ax = plt.subplots(1, 1, figsize=(14, 8))
    ax.axis('off')
    ax.set_facecolor(BG2_COLOR)
    ax.set_title('Эквивалентность Wyckoff ↔ SMC', color=TEXT_COLOR, fontsize=16, fontweight='bold', pad=20)

    # Table
    col_labels = ['Wyckoff', 'SMC', 'Рыночная логика']
    rows = [
        ['Spring / Shakeout', 'SSL Sweep', 'Снятие стопов покупателей\n— ложный пробой вниз'],
        ['Upthrust (UTAD)', 'BSL Sweep', 'Снятие стопов продавцов\n— ложный пробой вверх'],
        ['SOS (Sign of Strength)', 'Bullish BOS / MSS', 'Импульсный выход из\nдиапазона — сила быков'],
        ['LPS (Last Point of Support)', 'OB Retest / OTE', 'Последняя точка входа\nперед трендом'],
    ]
    table_data = [col_labels] + rows

    # Draw table manually with colored boxes
    col_widths = [0.25, 0.25, 0.40]
    x_start = 0.08
    y_start = 0.75
    row_height = 0.14

    for r, row in enumerate(table_data):
        y = y_start - r * row_height
        bg = '#111122' if r == 0 else None
        for c, (cell, w) in enumerate(zip(row, col_widths)):
            x = x_start + sum(col_widths[:c])
            clr = TEXT_COLOR if r == 0 else (ACCENT3 if c == 0 else (ACCENT2 if c == 1 else TEXT_DIM))
            if r == 0:
                ax.add_patch(Rectangle((x, y), w, row_height, facecolor='#333355', edgecolor='#555577'))
                ax.text(x + w/2, y + row_height/2, cell, fontsize=11, fontweight='bold',
                       color=TEXT_COLOR, ha='center', va='center', transform=ax.transAxes)
            else:
                ax.add_patch(Rectangle((x, y), w, row_height, facecolor=bg, edgecolor='#333355', linewidth=0.5))
                ax.text(x + 0.02, y + row_height/2, cell, fontsize=10, color=clr, va='center', transform=ax.transAxes)

    # Bottom diagram
    ax.text(0.5, 0.40, 'Синтез: Spring = SSL Sweep = Волна C Expanded Flat (Elliott)',
           fontsize=12, color=ACCENT1, ha='center', transform=ax.transAxes,
           bbox=dict(facecolor='#111122', alpha=0.8, edgecolor=ACCENT1, pad=8))
    ax.text(0.5, 0.25, 'Все три теории описывают одно и то же рыночное явление\nс разных ракурсов — финальную очистку перед истинным движением',
           fontsize=10, color=TEXT_DIM, ha='center', transform=ax.transAxes, style='italic')

    fig.tight_layout()
    return fig

def make_slide9_smc_structure():
    """Слайд 9: BOS vs CHoCH vs MSS"""
    data = load_1d('ADA')
    if not data: return None
    candles = get_candles(data)[-80:]
    if len(candles) < 20: return None

    closes = candle_array(candles, 'close')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    opens = candle_array(candles, 'open')
    dnums = dates_array(candles)

    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'SMC: BOS vs CHoCH vs MSS (ADA)', '')

    draw_candlestick(ax, candles)

    # Find swing points
    swings_h = []
    swings_l = []
    for i in range(2, len(candles)-2):
        if highs[i] > highs[i-1] and highs[i] > highs[i-2] and highs[i] > highs[i+1] and highs[i] > highs[i+2]:
            swings_h.append((dnums[i], highs[i]))
        if lows[i] < lows[i-1] and lows[i] < lows[i-2] and lows[i] < lows[i+1] and lows[i] < lows[i+2]:
            swings_l.append((dnums[i], lows[i]))

    # Draw swing lines
    for d, h in swings_h:
        ax.plot(d, h, '^', color=ACCENT1, markersize=6, alpha=0.8)
    for d, l in swings_l:
        ax.plot(d, l, 'v', color=ACCENT2, markersize=6, alpha=0.8)

    # Connect swings
    if len(swings_h) > 1:
        ds, vs = zip(*swings_h)
        ax.plot(ds, vs, color=ACCENT1, linewidth=0.8, linestyle='--', alpha=0.5)
    if len(swings_l) > 1:
        ds, vs = zip(*swings_l)
        ax.plot(ds, vs, color=ACCENT2, linewidth=0.8, linestyle='--', alpha=0.5)

    # Labels
    ax.text(0.02, 0.95, '▲ BOS = пробой структуры (продолжение)\n'
                        '▼ MSS = разворот с displacement > 1.5×ATR\n'
                        '◊ CHoCH = смена характера (предупреждение)',
           transform=ax.transAxes, fontsize=9, color=TEXT_DIM, va='top',
           bbox=dict(facecolor='#111122', alpha=0.7, edgecolor='none', pad=6))

    fig.tight_layout()
    return fig

def make_slide10_ob_fvg():
    """Слайд 10: Order Blocks + FVG"""
    data = load_1d('ADA')
    if not data: return None
    candles = get_candles(data)[-60:]
    if len(candles) < 10: return None

    closes = candle_array(candles, 'close')
    opens = candle_array(candles, 'open')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    dnums = dates_array(candles)

    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'Order Blocks и Fair Value Gaps (ADA)', '')

    dnums_fvg = matplotlib.dates.date2num(dates_array(candles))
    draw_candlestick(ax, candles)

    # Find FVGs
    for i in range(2, len(candles)):
        if lows[i-2] > highs[i]:  # Bullish FVG
            ax.fill_between([dnums_fvg[i-2], dnums_fvg[i]], highs[i], lows[i-2],
                           color=FVG_FILL, alpha=0.6)
            ax.annotate('FVG', ((dnums_fvg[i-2]+dnums_fvg[i])/2, (highs[i]+lows[i-2])/2),
                       fontsize=7, color=ACCENT2, ha='center',
                       bbox=dict(facecolor='#111122', alpha=0.7, edgecolor=ACCENT2, pad=2))

        if highs[i-2] < lows[i]:  # Bearish FVG
            ax.fill_between([dnums_fvg[i-2], dnums_fvg[i]], highs[i-2], lows[i],
                           color=FVG_FILL, alpha=0.6)
            ax.annotate('FVG', ((dnums_fvg[i-2]+dnums_fvg[i])/2, (highs[i-2]+lows[i])/2),
                       fontsize=7, color=ACCENT2, ha='center',
                       bbox=dict(facecolor='#111122', alpha=0.7, edgecolor=ACCENT2, pad=2))

    # Find potential Order Blocks (last opposite candle before impulse)
    for i in range(1, len(candles)):
        if closes[i] > opens[i] and closes[i-1] < opens[i-1]:
            # Bullish OB = last bear candle before green
            if i > 1 and abs(closes[i]-opens[i]) > abs(closes[i-1]-opens[i-1])*1.5:
                rect = Rectangle((dnums_fvg[i-1]-0.5, min(opens[i-1], closes[i-1])), 1,
                                abs(closes[i-1]-opens[i-1]),
                                facecolor=OB_BULL, edgecolor=BULL_COLOR, linewidth=1, linestyle='--')
                ax.add_patch(rect)
                ax.text(dnums_fvg[i-1], max(opens[i-1], closes[i-1]), ' OB ', fontsize=7,
                       color=BULL_COLOR, ha='center', va='bottom')

        if closes[i] < opens[i] and closes[i-1] > opens[i-1]:
            if i > 1 and abs(closes[i-1]-opens[i-1])*1.5 < abs(closes[i]-opens[i]):
                rect = Rectangle((dnums_fvg[i-1]-0.5, min(opens[i-1], closes[i-1])), 1,
                                abs(closes[i-1]-opens[i-1]),
                                facecolor=OB_BEAR, edgecolor=BEAR_COLOR, linewidth=1, linestyle='--')
                ax.add_patch(rect)
                ax.text(dnums_fvg[i-1], max(opens[i-1], closes[i-1]), ' OB ', fontsize=7,
                       color=BEAR_COLOR, ha='center', va='bottom')

    add_info_box(ax, 'OB: последняя противоположная свеча перед импульсом  |  FVG: Low[i-2] > High[i]')
    fig.tight_layout()
    return fig

def make_slide11_breaker_mitigation():
    """Слайд 11: Breaker vs Mitigation"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))
    ax1.axis('off'); ax2.axis('off')
    ax1.set_facecolor(BG2_COLOR); ax2.set_facecolor(BG2_COLOR)

    ax1.set_title('Breaker Block', color=ACCENT1, fontsize=14, fontweight='bold')
    ax2.set_title('Mitigation Block', color=ACCENT2, fontsize=14, fontweight='bold')

    # Breaker diagram
    # Draw schematic price action
    for ax_ref, title, is_breaker in [(ax1, 'Breaker Block', True), (ax2, 'Mitigation Block', False)]:
        # Schematic price path
        path = [(0.1, 0.7), (0.25, 0.85), (0.4, 0.6), (0.55, 0.9), (0.7, 0.5), (0.85, 0.8)]
        xs, ys = zip(*path)
        ax_ref.plot(xs, ys, color=TEXT_DIM, linewidth=1.5, transform=ax_ref.transAxes)

        # Sweep zone
        if is_breaker:
            ax_ref.axhspan(0.45, 0.55, xmin=0.45, xmax=0.65, facecolor=BEAR_COLOR, alpha=0.15, transform=ax_ref.transAxes)
            ax_ref.annotate('Sweep', (0.55, 0.5), fontsize=9, color=BEAR_COLOR, ha='center', transform=ax_ref.transAxes)

        # OB zone
        ob_y = 0.55 if is_breaker else 0.6
        ax_ref.axhspan(ob_y, ob_y+0.08, xmin=0.6, xmax=0.8, facecolor=OB_BULL, edgecolor=BULL_COLOR,
                      linestyle='--', linewidth=1, transform=ax_ref.transAxes)
        ax_ref.text(0.7, ob_y+0.04, 'OB', fontsize=10, color=BULL_COLOR, ha='center', transform=ax_ref.transAxes)

        # Description
        desc = ('Снятие ликвидности (Sweep)\n'
                '→ Стоп-хант пробойщиков\n'
                '→ Двойная ловушка\n'
                '→ Агрессивный отскок') if is_breaker else (
            'Отсутствие Sweep\n'
            '→ Failure Swing\n'
            '→ Смягчение (митигация)\n'
            '→ Упорядоченная коррекция')
        ax_ref.text(0.1, 0.15, desc, fontsize=9, color=TEXT_DIM, transform=ax_ref.transAxes)

    fig.suptitle('Breaker Block vs Mitigation Block', color=TEXT_COLOR, fontsize=15, fontweight='bold', y=0.98)
    fig.tight_layout()
    return fig

def make_slide12_crt_amd():
    """Слайд 12: CRT + AMD"""
    data = load_1d('ETH')
    if not data:
        data = load_1d('ADA')
    if not data: return None
    candles = get_candles(data)
    if len(candles) < 20: return None

    closes = candle_array(candles, 'close')
    dnums_raw = dates_array(candles)
    dnums_num = matplotlib.dates.date2num(dnums_raw)
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')

    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'CRT: Candle Range Theory — Цикл AMD (Accumulation-Manipulation-Distribution)', '')

    draw_candlestick(ax, candles)

    n = len(candles)
    # Divide into AMD zones
    zones = [
        (0, int(n*0.35), 'A — Accumulation', ACCENT1),
        (int(n*0.35), int(n*0.55), 'M — Manipulation (Sweep)', BEAR_COLOR),
        (int(n*0.55), n-1, 'D — Distribution', ACCENT2),
    ]
    for s, e, label, clr in zones:
        if e >= len(dnums_num): e = len(dnums_num)-1
        if s >= len(dnums_num): continue
        ax.axvspan(dnums_num[s]-0.5, dnums_num[e]+0.5, alpha=0.06, color=clr)
        mid = (s + e) // 2
        ax.text(dnums_num[mid], ax.get_ylim()[1], f'  {label}  ', fontsize=9, fontweight='bold',
               color=clr, ha='center', va='bottom',
               bbox=dict(facecolor='#111122', alpha=0.7, edgecolor=clr, pad=3))

    # Rule of 3 candles
    recent = candles[-10:]
    rd = dates_array(recent)
    ax.annotate('Правило 3 свечей:\nC1(импульс) → C2(FVG)\n→ C3(подтверждение)',
               xy=(rd[-3], closes[-3]), fontsize=8, color=ACCENT3,
               bbox=dict(facecolor='#111122', alpha=0.8, edgecolor=ACCENT3, pad=4))

    add_info_box(ax, 'AMD: Wyckoff → Williams → ICT → CRT (4-я итерация)')
    fig.tight_layout()
    return fig

def make_slide13_elliott():
    """Слайд 13: Elliott Wave + Expanded Flat"""
    data = load_1d('ATOM')
    if not data: return None
    candles = get_candles(data)[-150:]
    if len(candles) < 30: return None

    closes = candle_array(candles, 'close')
    dnums = dates_array(candles)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))
    style_ax(ax1, 'Elliott Wave: Импульс 5-3 (ATOM)', '')
    draw_candlestick(ax1, candles)

    # Mark wave structure (schematic)
    ax1.text(0.02, 0.95, 'Правила:\n1) 2 < 1\n2) 3 не самая короткая\n3) 4 не пересекает 1',
            transform=ax1.transAxes, fontsize=9, color=TEXT_DIM, va='top',
            bbox=dict(facecolor='#111122', alpha=0.7, edgecolor='none', pad=4))

    # Right: Expanded Flat diagram
    ax2.axis('off')
    ax2.set_title('Expanded Flat (3-3-5)', color=ACCENT1, fontsize=14, fontweight='bold')

    ax2.set_xlim(0, 10); ax2.set_ylim(0, 10)
    # Draw waves
    wave_points = [(1, 5), (3, 7), (4, 3), (6, 8.5), (7, 4), (9, 6)]
    labels = ['A', 'B', 'C', 'B\'', 'C\'', '']
    colors_wave = [BEAR_COLOR, BULL_COLOR, BEAR_COLOR, BULL_COLOR, BEAR_COLOR, TEXT_DIM]

    for i in range(len(wave_points)-1):
        x1, y1 = wave_points[i]
        x2, y2 = wave_points[i+1]
        ax2.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=colors_wave[i], lw=2.5))

    for (x, y), label, clr in zip(wave_points, labels, colors_wave):
        if label:
            ax2.text(x, y+0.3, label, fontsize=12, fontweight='bold', color=clr, ha='center')

    # Parameters
    ax2.text(5, 2, 'B = 105–138% от A (пробивает начало A)\nC = 162% от A (пробивает конец A)',
            fontsize=10, color=TEXT_DIM, ha='center',
            bbox=dict(facecolor='#111122', alpha=0.7, edgecolor=ACCENT1, pad=6))

    fig.tight_layout()
    return fig

def make_slide14_volume_profile():
    """Слайд 14: Volume Profile (POC, VA, HVN/LVN)"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)[-100:]
    if len(candles) < 10: return None

    closes = candle_array(candles, 'close')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    volumes = candle_array(candles, 'volume')
    dnums = dates_array(candles)

    # Calculate volume profile (horizontal histogram)
    price_min, price_max = np.min(lows), np.max(highs)
    step = (price_max - price_min) / 30
    if step == 0: step = 0.001
    price_bins = np.arange(price_min, price_max + step, step)
    vol_by_price = np.zeros(len(price_bins)-1)

    for i in range(len(candles)):
        for j in range(len(price_bins)-1):
            if price_bins[j] <= highs[i] and price_bins[j+1] >= lows[i]:
                vol_by_price[j] += volumes[i] / ((highs[i]-lows[i])/step + 1) if highs[i]-lows[i] > 0 else volumes[i]

    # Find POC
    poc_idx = np.argmax(vol_by_price)
    poc_price = (price_bins[poc_idx] + price_bins[poc_idx+1]) / 2

    # Value Area (68.2%)
    sorted_idx = np.argsort(-vol_by_price)
    cum_vol = 0
    total_vol = np.sum(vol_by_price)
    va_prices = []
    for idx in sorted_idx:
        cum_vol += vol_by_price[idx]
        va_prices.append((price_bins[idx] + price_bins[idx+1]) / 2)
        if cum_vol >= total_vol * 0.682:
            break
    va_min, va_max = min(va_prices), max(va_prices)

    fig = plt.figure(figsize=(16, 7))
    gs = fig.add_gridspec(1, 2, width_ratios=[3, 1], wspace=0.02)
    ax1 = fig.add_subplot(gs[0])
    ax2 = fig.add_subplot(gs[1])

    style_ax(ax1, 'Volume Profile: POC, Value Area, HVN/LVN (API3)', '')

    # Main chart
    draw_candlestick(ax1, candles)

    # POC line
    ax1.axhline(y=poc_price, color=ACCENT1, linewidth=1.5, linestyle='-', alpha=0.8, label=f'POC {price_format(poc_price)}')
    # Value Area
    ax1.axhspan(va_min, va_max, alpha=0.06, color=ACCENT2, label='Value Area (68.2%)')
    ax1.axhline(y=va_min, color=ACCENT2, linewidth=0.5, linestyle='--', alpha=0.5)
    ax1.axhline(y=va_max, color=ACCENT2, linewidth=0.5, linestyle='--', alpha=0.5)
    ax1.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # Profile histogram
    norm_vol = vol_by_price / np.max(vol_by_price) * 100
    ax2.barh(price_bins[:-1], norm_vol, height=step,
            color=VOLUME_COLOR, alpha=0.7)
    ax2.axhline(y=poc_price, color=ACCENT1, linewidth=1.5, label='POC')
    ax2.axhspan(va_min, va_max, alpha=0.1, color=ACCENT2)
    ax2.set_ylim(price_min, price_max)
    ax2.set_facecolor(BG2_COLOR)
    ax2.tick_params(colors=TEXT_DIM, labelleft=False)
    ax2.set_title('Объём', color=TEXT_DIM, fontsize=9)
    ax2.grid(False)

    add_info_box(ax1, f'POC: {price_format(poc_price)}  |  Value Area: {price_format(va_min)} – {price_format(va_max)}')
    fig.tight_layout()
    return fig

def make_slide15_triple_threat():
    """Слайд 15: Triple Threat"""
    fig, ax = plt.subplots(1, 1, figsize=(16, 8))
    ax.axis('off')
    ax.set_facecolor(BG2_COLOR)
    ax.set_title('Triple Threat — Слияние трёх факторов', color=TEXT_COLOR, fontsize=16, fontweight='bold', pad=20)

    # Three circles/boxes showing confluence
    boxes = [
        (0.08, 0.55, 'HVN (High Volume Node)', 'Узел максимального объёма\nЗона институционального\nинтереса', ACCENT1),
        (0.38, 0.55, 'Уровень S/R', 'Поддержка или сопротивление\nПодтверждённый историей\nуровень', '#7c4dff'),
        (0.68, 0.55, 'Anchored VWAP', 'Средневзвешенная цена\nот значимого экстремума\nДинамический уровень', ACCENT2),
    ]
    for x, y, title, desc, clr in boxes:
        ax.add_patch(FancyBboxPatch((x, y), 0.25, 0.30, boxstyle='round,pad=0.1',
                                     facecolor='#111122', edgecolor=clr, linewidth=2))
        ax.text(x+0.125, y+0.24, title, fontsize=11, fontweight='bold', color=clr, ha='center', transform=ax.transAxes)
        ax.text(x+0.125, y+0.10, desc, fontsize=9, color=TEXT_DIM, ha='center', transform=ax.transAxes)

    # Arrow
    ax.annotate('', xy=(0.33, 0.7), xytext=(0.25, 0.7),
               arrowprops=dict(arrowstyle='->', color=TEXT_DIM, lw=1.5), transform=ax.transAxes)
    ax.annotate('', xy=(0.63, 0.7), xytext=(0.55, 0.7),
               arrowprops=dict(arrowstyle='->', color=TEXT_DIM, lw=1.5), transform=ax.transAxes)

    # Result
    ax.add_patch(FancyBboxPatch((0.30, 0.10), 0.40, 0.25, boxstyle='round,pad=0.1',
                                 facecolor='#111122', edgecolor=ACCENT3, linewidth=3))
    ax.text(0.5, 0.27, '⚡ TRIPLE THREAT ⚡', fontsize=14, fontweight='bold', color=ACCENT3, ha='center', transform=ax.transAxes)
    ax.text(0.5, 0.15, 'Уровень высокого置信 — точка входа с максимальным совпадением',
           fontsize=9, color=TEXT_DIM, ha='center', transform=ax.transAxes)

    fig.tight_layout()
    return fig

def make_slide16_dom_depth():
    """Слайд 16: DOM Depth Chart"""
    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'DOM: Depth Chart — Глубина рынка Level 2', 'Объём (накопленный)', 'Цена')

    # Simulate DOM data
    center_price = 100.0
    prices_bid = np.linspace(center_price - 5, center_price, 50)
    prices_ask = np.linspace(center_price, center_price + 5, 50)
    bid_vol = np.exp(-np.linspace(0, 3, 50)) * 1000 + np.random.rand(50) * 100
    ask_vol = np.exp(-np.linspace(0, 3, 50)) * 1000 + np.random.rand(50) * 100
    bid_cum = np.cumsum(bid_vol[::-1])[::-1]
    ask_cum = np.cumsum(ask_vol)

    ax.fill_betweenx(prices_bid, 0, bid_cum, alpha=0.25, color=BULL_COLOR, label='Bids (лонги)')
    ax.plot(bid_cum, prices_bid, color=BULL_COLOR, linewidth=1.5)
    ax.fill_betweenx(prices_ask, 0, ask_cum, alpha=0.25, color=BEAR_COLOR, label='Asks (шорты)')
    ax.plot(ask_cum, prices_ask, color=BEAR_COLOR, linewidth=1.5)

    ax.axhline(y=center_price, color=ACCENT3, linewidth=2, linestyle='--', alpha=0.8, label='Spread')
    ax.legend(loc='upper right', fontsize=9, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    add_info_box(ax, 'Imbalance = (BidVol − AskVol) / (BidVol + AskVol)  |  Крутой slope = низкая ликвидность')
    fig.tight_layout()
    return fig

def make_slide17_of_imbalance():
    """Слайд 17: Order Flow Imbalance"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))

    # Left: Order Book Imbalance chart
    style_ax(ax1, 'Order Flow Imbalance (симуляция)', 'Время', 'Imbalance')
    np.random.seed(42)
    t = np.arange(100)
    imbalance = np.cumsum(np.random.randn(100) * 0.05)
    imbalance = np.clip(imbalance, -1, 1)
    ax1.fill_between(t, 0, imbalance, where=(imbalance>0), color=BULL_COLOR, alpha=0.3, label='Buy pressure')
    ax1.fill_between(t, 0, imbalance, where=(imbalance<0), color=BEAR_COLOR, alpha=0.3, label='Sell pressure')
    ax1.plot(t, imbalance, color=TEXT_COLOR, linewidth=1)
    ax1.axhline(y=0, color=TEXT_DIM, linewidth=0.5)
    ax1.axhline(y=0.3, color=BULL_COLOR, linestyle='--', alpha=0.5)
    ax1.axhline(y=-0.3, color=BEAR_COLOR, linestyle='--', alpha=0.5)
    ax1.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # Right: Metrics table
    ax2.axis('off')
    ax2.set_facecolor(BG2_COLOR)
    ax2.set_title('Online-индикаторы Order Flow', color=TEXT_COLOR, fontsize=13, fontweight='bold')

    metrics = [
        ('Real-Time CVD', 'Σ(Buy − Sell) per bar', 'Накопительная дельта'),
        ('Order Flow Imbalance', '(Buy−Sell)/(Buy+Sell)', 'Давление потока'),
        ('BAVR', 'Ask/(Ask+Bid)', 'Перекос книги'),
        ('Tape Speed', 'Принты/сек', 'Ритейл-активность'),
        ('DOM Heatmap', 'Bid/Ask по ценам', 'Концентрация'),
        ('Aggressive Ratio', 'Ask/Bid > 3.0', 'Агрессивный имбаланс'),
    ]
    y = 0.85
    for name, formula, desc in metrics:
        ax2.text(0.05, y, f'◆ {name}', fontsize=10, fontweight='bold', color=ACCENT2, transform=ax2.transAxes)
        ax2.text(0.05, y-0.04, f'  {formula}', fontsize=8, color=ACCENT3, transform=ax2.transAxes)
        ax2.text(0.05, y-0.07, f'  {desc}', fontsize=8, color=TEXT_DIM, transform=ax2.transAxes)
        y -= 0.13

    fig.tight_layout()
    return fig

def make_slide18_footprint():
    """Слайд 18: Footprint Stacked Imbalance"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))

    # Left: Cluster footprint grid
    ax1.axis('off')
    ax1.set_facecolor(BG2_COLOR)
    ax1.set_title('Footprint: Stacked Imbalance', color=TEXT_COLOR, fontsize=13, fontweight='bold')

    np.random.seed(123)
    grid_size = 15
    grid = np.random.rand(grid_size, grid_size) * 100
    # Create imbalance zone
    for i in range(5, 10):
        for j in range(6, 11):
            grid[i, j] = np.random.rand() * 50 + 150

    # Draw heatmap
    im = ax1.imshow(grid, cmap='RdYlGn', aspect='auto', interpolation='nearest')
    ax1.text(0.5, -0.05, '↑ Ask Volume', fontsize=9, color=BEAR_COLOR, ha='center', transform=ax1.transAxes)
    ax1.text(-0.05, 0.5, 'Price →', fontsize=9, color=TEXT_DIM, ha='center', rotation=90, transform=ax1.transAxes)

    # Highlight stacked imbalance
    ax1.add_patch(Rectangle((5.5, 5.5), 4, 4, fill=False, edgecolor=ACCENT3, linewidth=2, linestyle='--'))
    ax1.text(7.5, 4.5, 'Stacked Imbalance\n≥3 уровня > 300%', fontsize=8, color=ACCENT3, ha='center')

    # Right: Legend
    ax2.axis('off')
    ax2.set_facecolor(BG2_COLOR)
    footprint_info = [
        ('Stacked Imbalance', '≥3 уровня Ask/Bid > 300%\nСильная защита / агрессия', ACCENT3),
        ('Passive Absorption', 'Лимитный ордер поглощает\nагрессию → цена не движется', ACCENT2),
        ('Aggressive Imbalance', 'Ask(price)/Bid(price-1) > 3.0\nРыночный ордер доминирует', BEAR_COLOR),
        ('Delta', 'V(ask) − V(bid)\nРазница рыночных покупок/продаж', TEXT_COLOR),
        ('CVD Divergence', 'Цена ↑, CVD ↓\nСкрытая дистрибуция', ACCENT1),
    ]
    y = 0.9
    for name, desc, clr in footprint_info:
        ax2.text(0.1, y, f'◆ {name}', fontsize=10, fontweight='bold', color=clr, transform=ax2.transAxes)
        ax2.text(0.35, y-0.03, desc, fontsize=8, color=TEXT_DIM, transform=ax2.transAxes)
        y -= 0.16

    fig.tight_layout()
    return fig

def make_slide19_cvd():
    """Слайд 19: CVD дивергенция"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)[-60:]
    if len(candles) < 10: return None

    closes = candle_array(candles, 'close')
    volumes = candle_array(candles, 'volume')
    dnums = dates_array(candles)

    # Simulate CVD based on close changes
    delta = np.diff(closes, prepend=closes[0])
    cvd = np.cumsum(delta * np.sqrt(volumes) / 1000)

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(16, 7), gridspec_kw={'height_ratios': [2, 1]})
    style_ax(ax1, 'CVD: Cumulative Volume Delta — дивергенции (API3)', '')
    style_ax(ax2, 'CVD — подтверждение/опровержение тренда', '', '')

    draw_candlestick(ax1, candles)

    # CVD line
    ax2.plot(dnums, cvd, color=ACCENT2, linewidth=1.8, label='CVD')
    ax2.fill_between(dnums, 0, cvd, where=(cvd>0), color=BULL_COLOR, alpha=0.2)
    ax2.fill_between(dnums, 0, cvd, where=(cvd<0), color=BEAR_COLOR, alpha=0.2)
    ax2.axhline(y=0, color=TEXT_DIM, linewidth=0.5)
    ax2.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # Mark divergences (price up, CVD down)
    for i in range(5, len(candles)-5):
        if (closes[i] > closes[i-3] and cvd[i] < cvd[i-3] and
            cvd[i] < 0 and closes[i] > closes[i-5]):
            ax1.scatter(dnums[i], closes[i], color=ACCENT3, s=80, marker='o', facecolors='none', linewidth=2, zorder=5)
            ax1.annotate('CVD ↓\nPrice ↑', (dnums[i], closes[i]), fontsize=7, color=ACCENT3,
                        xytext=(10, 10), textcoords='offset points',
                        arrowprops=dict(arrowstyle='->', color=ACCENT3, lw=0.5))

    add_info_box(ax1, 'Медвежья дивергенция CVD: цена растёт → CVD падает = скрытые продажи')
    fig.tight_layout()
    return fig

def make_slide20_iceberg():
    """Слайд 20: Iceberg Detection"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))

    # Left: Iceberg schematic
    ax1.axis('off')
    ax1.set_facecolor(BG2_COLOR)
    ax1.set_title('Iceberg Order Detection', color=TEXT_COLOR, fontsize=13, fontweight='bold')

    # Draw iceberg visual
    ax1.add_patch(Rectangle((0.2, 0.3), 0.6, 0.15, facecolor=ACCENT2, alpha=0.2, edgecolor=ACCENT2, linewidth=1))
    ax1.text(0.5, 0.37, 'Visible DOM: 500 BTC', fontsize=10, color=ACCENT2, ha='center', transform=ax1.transAxes)

    ax1.add_patch(Rectangle((0.25, 0.1), 0.5, 0.20, facecolor=ACCENT2, alpha=0.4, edgecolor=ACCENT2, linewidth=1, linestyle='--'))
    ax1.text(0.5, 0.20, 'Hidden Reserve', fontsize=10, color=TEXT_DIM, ha='center', transform=ax1.transAxes)
    ax1.text(0.5, 0.05, 'T&S: 8500 BTC за 30 сек', fontsize=9, color=ACCENT3, ha='center', transform=ax1.transAxes)

    ax1.annotate('', xy=(0.8, 0.37), xytext=(0.9, 0.55),
                arrowprops=dict(arrowstyle='->', color=BEAR_COLOR, lw=2), transform=ax1.transAxes)
    ax1.text(0.85, 0.58, 'T&S Vol > 5× DOM Vol\nСтатичный DOM размер', fontsize=8, color=BEAR_COLOR,
            ha='center', transform=ax1.transAxes)

    # Right: real example (NVDA)
    ax2.axis('off')
    ax2.set_facecolor(BG2_COLOR)
    ax2.set_title('Кейс: NVDA $875.00', color=ACCENT1, fontsize=13, fontweight='bold')

    case_text = (
        'Уровень: $875.00\n\n'
        'DOM: Ask = 500 акций\n'
        'Лента: 12 покупок\n'
        '  по 500–2000 акций\n'
        '  за 30 секунд\n\n'
        'После каждой сделки\n'
        'Ask обновляется до 500\n\n'
        'Вывод: Айсберг —\n'
        'скрытая дистрибуция\n'
        'в сильный спрос'
    )
    ax2.text(0.1, 0.85, case_text, fontsize=10, color=TEXT_DIM, transform=ax2.transAxes, verticalalignment='top')

    fig.suptitle('Iceberg Orders — Скрытая ликвидность', color=TEXT_COLOR, fontsize=15, fontweight='bold', y=0.98)
    fig.tight_layout()
    return fig

def make_slide21_ote():
    """Слайд 21: OTE / Premium-Discount"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)[-50:]
    if len(candles) < 5: return None

    closes = candle_array(candles, 'close')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    dnums = dates_array(candles)

    # Find recent swing high/low
    swing_high = np.max(highs[-20:])
    swing_low = np.min(lows[-20:])
    rng = swing_high - swing_low
    fib_62 = swing_low + rng * (1 - 0.62)
    fib_79 = swing_low + rng * (1 - 0.79)
    sweet_spot = swing_low + rng * (1 - 0.705)

    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'OTE: Optimal Trade Entry — Premium & Discount Zones (API3)', '')

    draw_candlestick(ax, candles)

    # Premium/Discount zones
    mid = (swing_high + swing_low) / 2
    ax.axhspan(mid, swing_high, alpha=0.06, color=BEAR_COLOR, label='Premium Zone')
    ax.axhspan(swing_low, mid, alpha=0.06, color=BULL_COLOR, label='Discount Zone')

    # Fib levels
    fib_levels = [(0.0, swing_low), (0.236, swing_low + rng*0.236), (0.382, swing_low + rng*0.382),
                  (0.5, mid), (0.618, swing_low + rng*0.618), (0.786, swing_low + rng*0.786), (1.0, swing_high)]
    for label, price in fib_levels:
        ax.axhline(y=price, color=TEXT_DIM, linewidth=0.5, alpha=0.5)
        ax.text(dnums[-1], price, f'  {label}', fontsize=8, color=TEXT_DIM)

    # OTE zone
    ax.axhspan(fib_62, fib_79, alpha=0.15, color=ACCENT2, label='OTE Zone (62–79%)')
    ax.axhline(y=sweet_spot, color=ACCENT3, linewidth=1.5, linestyle='--', label=f'Sweet Spot {price_format(sweet_spot)}')

    ax.axhline(y=swing_low, color=TEXT_DIM, linewidth=0.8, alpha=0.6)
    ax.axhline(y=swing_high, color=TEXT_DIM, linewidth=0.8, alpha=0.6)

    ax.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    add_info_box(ax, f'OTE: {price_format(fib_79)} – {price_format(fib_62)}  |  Sweet Spot: {price_format(sweet_spot)}')
    fig.tight_layout()
    return fig

def make_slide22_heatmap():
    """Слайд 22: Liquidation Heatmap"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)[-100:]
    if len(candles) < 10: return None

    closes = candle_array(candles, 'close')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    dnums = dates_array(candles)

    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'Liquidation Heatmap — Магнитные зоны ликвидаций (API3)', '')

    draw_candlestick(ax, candles)

    # Simulate liquidation clusters
    np.random.seed(42)
    current_price = closes[-1]
    for i in range(10):
        liq_price = current_price * (1 + np.random.randn() * 0.02)
        liq_time = np.random.choice(dnums)
        liq_size = np.random.rand() * 100
        alpha = min(liq_size / 50, 0.8)
        ax.scatter(liq_time, liq_price, s=liq_size*5, c=ACCENT3, alpha=alpha, edgecolors='none', zorder=3)

    # Entry price line from card if avail
    card = load_card('API3')
    if card:
        data = card.get('data', {})
        live = data.get('live_position', {})
        if live.get('hold_side') == 'long':
            liq = live.get('liquidation_price')
            entry = data.get('emoji_entry', {}).get('entry_price')
            if liq and entry:
                try:
                    liq_f = float(liq)
                    entry_f = float(entry)
                    ax.axhline(y=entry_f, color=BULL_COLOR, linewidth=1.5, linestyle='-', alpha=0.8,
                              label=f'Entry: {price_format(entry_f)}')
                    ax.axhline(y=liq_f, color=BEAR_COLOR, linewidth=1.5, linestyle='--', alpha=0.8,
                              label=f'Lp: {price_format(liq_f)}')
                    ax.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)
                except:
                    pass

    add_info_box(ax, 'Lp(LONG) = Entry × (1 − 1/Leverage)  |  Яркая зона = магнит для цены')
    fig.tight_layout()
    return fig

def make_slide23_oi():
    """Слайд 23: OI Trend Health (4 сценария)"""
    fig, axes = plt.subplots(2, 2, figsize=(16, 8))
    axes = axes.flatten()

    scenarios = [
        ('Здоровый тренд', [1, 2, 3, 4, 5], [1, 2, 3, 4, 5], [0.01]*5, BULL_COLOR,
         'Price ↑, OI ↑\nFunding ~0.01%\nПриток капитала'),
        ('Leverage Flush', [1, 2, 4, 6, 8], [1, 2, 3, 4, 5], [0.02, 0.03, 0.05, 0.08, 0.12], BEAR_COLOR,
         'Price ↑↑, OI ↑\nFunding > 0.05%\nРиск каскада'),
        ('Затухание', [1, 2, 3, 4, 5], [5, 4.5, 4, 3.5, 3], [0.01]*5, ACCENT1,
         'Price ↑, OI ↓\nФиксация китами\nВыход'),
        ('Капитуляция', [5, 4, 3, 2, 1], [5, 3.5, 2, 1, 0.5], [-0.01, -0.02, -0.03, -0.04, -0.05], '#7c4dff',
         'Price ↓↓, OI ↓↓\nМассовые ликвидации\nПоиск дна'),
    ]

    for idx, (title, price, oi, funding, clr, desc) in enumerate(scenarios):
        ax = axes[idx]
        ax.set_facecolor(BG2_COLOR)
        ax.set_title(title, color=clr, fontsize=11, fontweight='bold')

        t = np.arange(len(price))
        ax.plot(t, price, color=TEXT_COLOR, linewidth=2, label='Price')
        ax2 = ax.twinx()
        ax2.plot(t, oi, color=ACCENT2, linewidth=2, linestyle='--', label='OI')
        ax2.axhline(y=0, color=TEXT_DIM, linewidth=0.3)

        # Funding markers
        for i, f in enumerate(funding):
            if abs(f) > 0.05:
                ax.scatter(i, price[i], color=ACCENT3, s=50, zorder=5, marker='^')
            elif abs(f) > 0.02:
                ax.scatter(i, price[i], color=TEXT_DIM, s=30, zorder=5, marker='o', alpha=0.5)

        ax.tick_params(colors=TEXT_DIM, labelsize=7)
        ax2.tick_params(colors=TEXT_DIM, labelsize=7)
        ax.text(0.5, -0.25, desc, fontsize=8, color=TEXT_DIM, ha='center', transform=ax.transAxes)

    fig.suptitle('OI Trend Health — 4 сценария', color=TEXT_COLOR, fontsize=14, fontweight='bold')
    fig.tight_layout()
    return fig

def make_slide24_sentiment():
    """Слайд 24: Fear & Greed + Funding"""
    fig, axes = plt.subplots(1, 3, figsize=(16, 7))

    # Fear & Greed gauge
    ax1 = axes[0]
    ax1.axis('off')
    ax1.set_facecolor(BG2_COLOR)
    ax1.set_title('Fear & Greed Index', color=TEXT_COLOR, fontsize=11, fontweight='bold')

    # Draw gauge manually
    gauge_colors = ['#ff1744', '#ff9100', '#ffea00', '#00c853']
    for i, (label, clr) in enumerate([('Fear', BEAR_COLOR), ('Neutral', ACCENT1), ('Greed', BULL_COLOR)]):
        ax1.add_patch(Rectangle((0.1 + i*0.26, 0.4), 0.24, 0.08, facecolor=clr, alpha=0.3, edgecolor=clr))
        ax1.text(0.22 + i*0.26, 0.36, label, fontsize=8, color=TEXT_DIM, ha='center', transform=ax1.transAxes)

    # Needle at Greed=75 (current)
    needle = 0.75
    x_pos = 0.1 + needle * 0.78
    ax1.plot(x_pos, 0.44, 'v', color=ACCENT3, markersize=12)
    ax1.text(x_pos, 0.50, '75\n(Extreme\nGreed)', fontsize=9, color=ACCENT3, ha='center', transform=ax1.transAxes,
            bbox=dict(facecolor='#111122', alpha=0.7, edgecolor=ACCENT3, pad=3))

    ax1.text(0.5, 0.15, 'Сигнал: Extreme Greed\n→ риск вершины\n→ искать шорты', fontsize=9, color=TEXT_DIM,
            ha='center', transform=ax1.transAxes, bbox=dict(facecolor='#111122', alpha=0.5, edgecolor='none', pad=4))

    # Funding Rate gauge
    ax2 = axes[1]
    ax2.axis('off')
    ax2.set_facecolor(BG2_COLOR)
    ax2.set_title('Funding Rate Sentiment', color=TEXT_COLOR, fontsize=11, fontweight='bold')

    for i, (label, clr) in enumerate([('Neg\nShorts pay', BULL_COLOR), ('Neutral', TEXT_DIM), ('Pos\nLongs pay', BEAR_COLOR)]):
        ax2.add_patch(Rectangle((0.1 + i*0.26, 0.4), 0.24, 0.08, facecolor=clr, alpha=0.3, edgecolor=clr))
        ax2.text(0.22 + i*0.26, 0.35, label, fontsize=7, color=TEXT_DIM, ha='center', transform=ax2.transAxes)

    # Needle at 0.08% (positive -> risk)
    x_pos = 0.1 + 0.65 * 0.78  # 65% of positive zone
    ax2.plot(x_pos, 0.44, 'v', color=BEAR_COLOR, markersize=12)
    ax2.text(x_pos, 0.50, '+0.08%\nLeverage Flush\nRisk!', fontsize=8, color=BEAR_COLOR, ha='center', transform=ax2.transAxes,
            bbox=dict(facecolor='#111122', alpha=0.7, edgecolor=BEAR_COLOR, pad=3))

    # Put/Call
    ax3 = axes[2]
    ax3.axis('off')
    ax3.set_facecolor(BG2_COLOR)
    ax3.set_title('Дополнительные метрики', color=TEXT_COLOR, fontsize=11, fontweight='bold')

    metrics_sent = [
        ('Put/Call Ratio', '> 1.0 = страх\n< 0.5 = жадность'),
        ('VIX / RVX', '< 15 = спокойствие\n> 30 = паника'),
        ('Social Volume', 'Рост = хайп\n(контрсигнал)'),
        ('COT Report', 'Commercial vs\nSpeculator позиции'),
    ]
    y = 0.85
    for name, desc in metrics_sent:
        ax3.text(0.05, y, f'◆ {name}', fontsize=9, fontweight='bold', color=ACCENT2, transform=ax3.transAxes)
        ax3.text(0.3, y-0.02, desc, fontsize=8, color=TEXT_DIM, transform=ax3.transAxes)
        y -= 0.18

    fig.suptitle('Sentiment: Настроение рынка', color=TEXT_COLOR, fontsize=14, fontweight='bold')
    fig.tight_layout()
    return fig

def make_slide25_onchain():
    """Слайд 25: On-Chain: MVRV / SOPR"""
    fig, axes = plt.subplots(2, 2, figsize=(16, 8))
    axes = axes.flatten()

    # Synthetic on-chain data
    np.random.seed(42)
    days = 365
    t = np.arange(days)

    # MVRV
    ax = axes[0]
    style_ax(ax, 'MVRV Ratio', 'Days', 'MVRV')
    mvrv = 1 + 0.5 * np.sin(t / 100) + np.random.randn(days) * 0.05
    ax.plot(t, mvrv, color=ACCENT2, linewidth=1.5)
    ax.axhline(y=3.5, color=BEAR_COLOR, linestyle='--', alpha=0.7, label='Overheat > 3.5')
    ax.axhline(y=1.0, color=BULL_COLOR, linestyle='--', alpha=0.7, label='Fair = 1.0')
    ax.axhline(y=0.5, color=ACCENT1, linestyle='--', alpha=0.7, label='Capitulation < 0.5')
    ax.fill_between(t, 0, mvrv, where=(mvrv<0.5), color=BULL_COLOR, alpha=0.1)
    ax.fill_between(t, 0, mvrv, where=(mvrv>3.5), color=BEAR_COLOR, alpha=0.1)
    ax.legend(fontsize=7, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # SOPR
    ax = axes[1]
    style_ax(ax, 'SOPR', 'Days', 'SOPR')
    sopr = 1 + 0.1 * np.sin(t / 50) + np.random.randn(days) * 0.03
    ax.plot(t, sopr, color='#7c4dff', linewidth=1.5)
    ax.axhline(y=1.0, color=TEXT_COLOR, linestyle='--', alpha=0.5, label='Break-even')
    ax.fill_between(t, 0, sopr, where=(sopr<1), color=BEAR_COLOR, alpha=0.1, label='Loss')
    ax.legend(fontsize=7, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # NVT
    ax = axes[2]
    style_ax(ax, 'NVT Ratio', 'Days', 'NVT')
    nvt = 80 + 30 * np.sin(t / 70) + np.random.randn(days) * 5
    ax.plot(t, nvt, color=ACCENT1, linewidth=1.5)
    ax.axhline(y=100, color=BEAR_COLOR, linestyle='--', alpha=0.5, label='Overvalued > 100')
    ax.axhline(y=50, color=BULL_COLOR, linestyle='--', alpha=0.5, label='Undervalued < 50')
    ax.legend(fontsize=7, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # Summary
    ax = axes[3]
    ax.axis('off')
    ax.set_facecolor(BG2_COLOR)
    ax.set_title('On-Chain: Сигналы', color=TEXT_COLOR, fontsize=12, fontweight='bold')

    signals = [
        ('MVRV > 3.5', 'Перегрев, вершина', BEAR_COLOR),
        ('MVRV < 1.0', 'Дно, накопление', BULL_COLOR),
        ('SOPR < 1.0', 'Капитуляция', '#7c4dff'),
        ('NVT > 100', 'Переоценен', BEAR_COLOR),
        ('NVT < 50', 'Недооценен', BULL_COLOR),
        ('Puell < 0.5', 'Дно майнеров (бычий)', ACCENT2),
        ('Whale Inflow ↑', 'Киты несут на биржи', BEAR_COLOR),
        ('Stablecoin Inflow ↑', 'Покупки на подходе', BULL_COLOR),
    ]
    y = 0.85
    for name, desc, clr in signals:
        ax.text(0.05, y, f'{name}', fontsize=8, color=clr, fontweight='bold', transform=ax.transAxes)
        ax.text(0.5, y, f'→ {desc}', fontsize=8, color=TEXT_DIM, transform=ax.transAxes)
        y -= 0.10

    fig.suptitle('On-Chain метрики: MVRV, SOPR, NVT', color=TEXT_COLOR, fontsize=14, fontweight='bold')
    fig.tight_layout()
    return fig

def make_slide26_whales():
    """Слайд 26: Whale Activity"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 7))

    # Exchange flows
    style_ax(ax1, 'Exchange Netflow — Движение китов', 'Days', 'BTC')
    np.random.seed(7)
    days = 100
    t = np.arange(days)
    netflow = np.cumsum(np.random.randn(days) * 50)
    ax1.fill_between(t, 0, netflow, where=(netflow>0), color=BEAR_COLOR, alpha=0.2, label='Inflow (bearish)')
    ax1.fill_between(t, 0, netflow, where=(netflow<0), color=BULL_COLOR, alpha=0.2, label='Outflow (bullish)')
    ax1.plot(t, netflow, color=TEXT_COLOR, linewidth=1.5)
    ax1.axhline(y=0, color=TEXT_DIM, linewidth=0.5)
    ax1.legend(fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    # Metrics list
    ax2.axis('off')
    ax2.set_facecolor(BG2_COLOR)
    ax2.set_title('Whale & Exchange Metrics', color=TEXT_COLOR, fontsize=12, fontweight='bold')

    whale_metrics = [
        ('Whale Tx Count', 'Транзакции > $100k / $1M / $10M'),
        ('Whale Concentration', '% supply в топ-10 адресах'),
        ('Exchange Whale Ratio', 'Топ-10 inflow / total inflow'),
        ('Exchange Balance', 'Резервы на биржах'),
        ('Stablecoin Inflow', 'USDT/USDC на биржи'),
        ('BTC Reserve Risk', 'Долгосрочный индикатор'),
        ('Active Addresses', 'Уникальные адреса/день'),
        ('CDD', 'Coin Days Destroyed'),
    ]
    y = 0.88
    for name, desc in whale_metrics:
        ax2.text(0.05, y, f'◆ {name}', fontsize=9, fontweight='bold', color=ACCENT2, transform=ax2.transAxes)
        ax2.text(0.3, y-0.02, desc, fontsize=8, color=TEXT_DIM, transform=ax2.transAxes)
        y -= 0.11

    fig.suptitle('Киты и биржевые потоки: On-Chain анализ', color=TEXT_COLOR, fontsize=14, fontweight='bold')
    fig.tight_layout()
    return fig

def make_slide27_checklist():
    """Слайд 27: 10-Level Checklist"""
    fig, ax = plt.subplots(1, 1, figsize=(16, 8))
    ax.axis('off')
    ax.set_facecolor(BG2_COLOR)
    ax.set_title('10-уровневый чек-лист входа', color=TEXT_COLOR, fontsize=16, fontweight='bold', pad=20)

    levels = [
        (1, 'Wyckoff Phase', 'Макро-фаза: накопление/распределение', ACCENT1),
        (2, 'MSS', 'Слом структуры с displacement > 1.5×ATR', BULL_COLOR),
        (3, 'Sweep', 'Снятие ликвидности BSL/SSL', BEAR_COLOR),
        (4, 'Value Area', 'VAH для шорта, VAL для лонга', ACCENT2),
        (5, 'POI', 'Order Block или FVG на пути', '#7c4dff'),
        (6, 'OTE', 'Fib 70.5% (Sweet Spot)', ACCENT3),
        (7, 'Heatmap', 'Кластер ликвидаций поблизости', ACCENT1),
        (8, 'OI Health', 'Trend Health = цена+OI+funding', ACCENT2),
        (9, 'Footprint', 'Stacked Imbalance / Absorption', BULL_COLOR),
        (10, 'RR > 1:2', 'Risk/Reward не менее 1 к 2', BEAR_COLOR),
    ]

    for i, (num, name, desc, clr) in enumerate(levels):
        y = 0.85 - i * 0.075
        # Number badge
        ax.add_patch(Rectangle((0.05, y-0.02), 0.04, 0.04, facecolor=clr, alpha=0.3, edgecolor=clr, linewidth=1))
        ax.text(0.07, y, str(num), fontsize=9, fontweight='bold', color=clr, ha='center', va='center', transform=ax.transAxes)
        # Name
        ax.text(0.12, y, name, fontsize=10, fontweight='bold', color=clr, transform=ax.transAxes)
        # Description
        ax.text(0.35, y, desc, fontsize=9, color=TEXT_DIM, transform=ax.transAxes)

    # Status box
    ax.text(0.5, 0.05, 'Сеттап считается валидным при min 4 из 6 ключевых подтверждений',
           fontsize=12, color=ACCENT3, ha='center', transform=ax.transAxes,
           bbox=dict(facecolor='#111122', alpha=0.8, edgecolor=ACCENT3, pad=8))

    fig.tight_layout()
    return fig

def make_slide28_velvet():
    """Слайд 28: Кейс VELVETUSDT"""
    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'Кейс: VELVETUSDT — Pump & Dump (июнь 2026)', '')

    # Simulate the 17-minute pump & dump
    np.random.seed(17)
    n = 60  # 60 data points for 17 min
    t = np.linspace(0, 17, n)
    base = 1.3248
    # Pump: +5.4% in 8 min, then dump back in 9 min
    pump = np.exp(-((t-4)/3.5)**2) * 0.054 * base
    price = base + pump
    # Add noise
    price += np.random.randn(n) * 0.002
    # Make it look like candles
    open_p = price[:-1]
    close_p = price[1:]
    high_p = np.maximum(open_p, close_p) + np.random.rand(n-1)*0.005
    low_p = np.minimum(open_p, close_p) - np.random.rand(n-1)*0.005

    for i in range(n-1):
        color = BULL_COLOR if close_p[i] >= open_p[i] else BEAR_COLOR
        ax.vlines(i, low_p[i], high_p[i], color=color, linewidth=0.5, alpha=0.5)
        rect = Rectangle((i-0.3, min(open_p[i], close_p[i])), 0.6, abs(close_p[i]-open_p[i]),
                        facecolor=color, edgecolor=color, alpha=0.8)
        ax.add_patch(rect)

    # Annotations
    ax.axhline(y=base, color=TEXT_DIM, linestyle='--', alpha=0.6, label=f'Origin: {base}')
    peak_idx = np.argmax(close_p)
    ax.annotate(f'Peak: {close_p[peak_idx]:.4f}\n+5.4%', xy=(peak_idx, close_p[peak_idx]),
               fontsize=9, color=ACCENT3,
               arrowprops=dict(arrowstyle='->', color=ACCENT3, lw=1))

    end_idx = n-2
    ax.annotate(f'Return: {close_p[end_idx]:.4f}\n17 min 9 sec', xy=(end_idx, close_p[end_idx]),
               fontsize=9, color=BEAR_COLOR,
               arrowprops=dict(arrowstyle='->', color=BEAR_COLOR, lw=1))

    ax.set_xlabel('Время (минуты)', color=TEXT_DIM)
    ax.set_ylabel('Цена', color=TEXT_DIM)
    ax.legend(fontsize=9, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    add_info_box(ax, 'Институциональное предложение абсорбировало ликвидность ритейла → Full Return')
    fig.tight_layout()
    return fig

def make_slide29_api3():
    """Слайд 29: Кейс API3 #4"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)
    if len(candles) < 20: return None

    closes = candle_array(candles, 'close')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    volumes = candle_array(candles, 'volume')
    dnums = dates_array(candles)

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(16, 8), gridspec_kw={'height_ratios': [3, 1]})
    style_ax(ax1, 'Кейс: API3 #4 — 195 дней в позиции (реальные данные)', '')
    style_ax(ax2, 'Объём', '', 'Volume')

    draw_candlestick(ax1, candles)

    # Entry and liquidation from card
    card = load_card('API3')
    if card:
        d = card.get('data', {})
        live = d.get('live_position', {})
        entry_data = d.get('emoji_entry', {})
        entry_price = float(entry_data.get('entry_price', 0.2992))
        liq_price = float(live.get('liquidation_price', -0.663))
        leverage = float(live.get('leverage', 10))

        ax1.axhline(y=entry_price, color=BULL_COLOR, linewidth=1.5, linestyle='-',
                   alpha=0.8, label=f'Entry: {entry_price}')
        if liq_price < 0:
            liq_price = entry_price * (1 - 1/leverage)
        ax1.axhline(y=liq_price, color=BEAR_COLOR, linewidth=1.5, linestyle='--',
                   alpha=0.8, label=f'Lp (10x): ~{liq_price:.4f}')
        ax1.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

        ax1.annotate(f'Entry\n{entry_price}', xy=(dnums[0], entry_price),
                    fontsize=8, color=BULL_COLOR,
                    arrowprops=dict(arrowstyle='->', color=BULL_COLOR, lw=1))

    # Volume
    ax2.bar(dnums, volumes, color=VOLUME_COLOR, width=0.6)
    ax2.set_yscale('log')

    add_info_box(ax1, 'Параметры: Entry 0.2992, 10x, Liquidation ~0.2693  |  Длительность: 195 дней, ROE: -329%')
    fig.tight_layout()
    return fig

def make_slide30_btc():
    """Слайд 30: BTC Creek & Ice"""
    data = load_1d('API3')
    if not data: return None
    candles = get_candles(data)
    if len(candles) < 50: return None

    closes = candle_array(candles, 'close')
    highs = candle_array(candles, 'high')
    lows = candle_array(candles, 'low')
    dnums = dates_array(candles)

    fig, ax = plt.subplots(1, 1, figsize=(16, 7))
    style_ax(ax, 'BTC: Creek & Ice — Институциональный диапазон 2026', '')

    draw_candlestick(ax, candles)

    # Find Creek (resistance) and Ice (support) zones
    recent = candles[-40:]
    r_high = max(c['high'] for c in recent)
    r_low = min(c['low'] for c in recent)

    # Equilibrium
    eq = (r_high + r_low) / 2

    ax.axhline(y=r_high, color=BEAR_COLOR, linewidth=1.5, linestyle='--', alpha=0.7, label=f'Creek: {r_high:.4f}')
    ax.axhline(y=r_low, color=BULL_COLOR, linewidth=1.5, linestyle='--', alpha=0.7, label=f'Ice: {r_low:.4f}')
    ax.axhline(y=eq, color=ACCENT1, linewidth=1.5, linestyle='-', alpha=0.8, label=f'Equilibrium: {eq:.4f}')

    # POC from earlier
    volumes = candle_array(recent, 'volume')
    v_max_idx = np.argmax(volumes)
    poc_price = (recent[v_max_idx]['high'] + recent[v_max_idx]['low']) / 2
    ax.axhline(y=poc_price, color=ACCENT2, linewidth=1, linestyle=':', alpha=0.6, label=f'POC: {poc_price:.4f}')

    ax.legend(loc='upper right', fontsize=8, facecolor='#111122', edgecolor='#444466', labelcolor=TEXT_COLOR)

    annotations_text = (
        'Сценарий:\n'
        '1. Снятие ликвидности за Ice\n'
        '2. Возврат к Equilibrium\n'
        '3. POC = зона накопления'
    )
    ax.text(0.02, 0.95, annotations_text, fontsize=9, color=TEXT_DIM, va='top', transform=ax.transAxes,
           bbox=dict(facecolor='#111122', alpha=0.7, edgecolor='none', pad=6))

    add_info_box(ax, f'Creek={r_high:.4f}  |  Ice={r_low:.4f}  |  Eq={eq:.4f}  |  POC={poc_price:.4f}', y=0.02, fontsize=8)
    fig.tight_layout()
    return fig

# ─── MAIN ─────────────────────────────────────────────
def main():
    print("=" * 60)
    print("Генерация презентации учебника институционального трейдинга")
    print("=" * 60)

    # Generate all slide figures
    generators = [
        ('ADX-фильтр', make_slide1_adx),
        ('MACD + Stochastic', make_slide2_macd_stoch),
        ('Bollinger + ATR', make_slide3_bb_atr),
        ('Свечные паттерны', make_slide4_candle_patterns),
        ('Тройные паттерны', make_slide5_triple_star),
        ('Графические паттерны', make_slide6_graphical_patterns),
        ('Wyckoff A–E', make_slide7_wyckoff),
        ('Wyckoff ↔ SMC', make_slide8_wyckoff_smc),
        ('BOS/CHoCH/MSS', make_slide9_smc_structure),
        ('OB + FVG', make_slide10_ob_fvg),
        ('Breaker vs Mitigation', make_slide11_breaker_mitigation),
        ('CRT + AMD', make_slide12_crt_amd),
        ('Elliott + Expanded Flat', make_slide13_elliott),
        ('Volume Profile', make_slide14_volume_profile),
        ('Triple Threat', make_slide15_triple_threat),
        ('DOM Depth Chart', make_slide16_dom_depth),
        ('Order Flow Imbalance', make_slide17_of_imbalance),
        ('Footprint', make_slide18_footprint),
        ('CVD', make_slide19_cvd),
        ('Iceberg Detection', make_slide20_iceberg),
        ('OTE Zone', make_slide21_ote),
        ('Liquidation Heatmap', make_slide22_heatmap),
        ('OI Trend Health', make_slide23_oi),
        ('Sentiment', make_slide24_sentiment),
        ('On-Chain', make_slide25_onchain),
        ('Whale Activity', make_slide26_whales),
        ('10-Level Checklist', make_slide27_checklist),
        ('Кейс VELVETUSDT', make_slide28_velvet),
        ('Кейс API3 #4', make_slide29_api3),
        ('Кейс BTC Creek & Ice', make_slide30_btc),
    ]

    slides = []
    for name, gen in generators:
        print(f"  → Слайд: {name}...", end=' ')
        try:
            fig = gen()
            if fig is not None:
                slides.append((name, fig))
                print(f"✅ ({fig.get_size_inches()})")
            else:
                print(f"⚠️  нет данных, пропущен")
        except Exception as e:
            print(f"❌ {e}")
            import traceback; traceback.print_exc()

    print(f"\n✅ Сгенерировано слайдов: {len(slides)}/{len(generators)}")

    # ─── Build PPTX ────────────────────────────────────
    print("\nСборка PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    from pptx.util import Pt as PtSize
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Учебник институционального трейдинга"
    p.font.size = PtSize(36)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "От Price Action до рыночной микроструктуры"
    p2.font.size = PtSize(18)
    p2.font.color.rgb = RGBColor(0xff, 0xea, 0x00)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "30 визуализаций на основе реальных данных проекта Fundament RF"
    p3.font.size = PtSize(12)
    p3.font.color.rgb = RGBColor(0xcc, 0xcc, 0xdd)
    p3.alignment = PP_ALIGN.CENTER

    # ─── Build PDF (via matplotlib PdfPages) ────────────
    from matplotlib.backends.backend_pdf import PdfPages
    pdf_path = OUTPUT_DIR / 'Презентация_учебника.pdf'
    print(f"\nСборка PDF...")
    with PdfPages(str(pdf_path)) as pdf:
        # Title page
        fig_t = plt.figure(figsize=(16, 9))
        fig_t.patch.set_facecolor(BG_COLOR)
        ax_t = fig_t.add_subplot(111)
        ax_t.axis('off')
        ax_t.set_facecolor(BG_COLOR)
        ax_t.text(0.5, 0.6, 'Учебник институционального трейдинга', fontsize=32, fontweight='bold',
                 color=TEXT_COLOR, ha='center', transform=ax_t.transAxes)
        ax_t.text(0.5, 0.5, 'От Price Action до рыночной микроструктуры', fontsize=18,
                 color=ACCENT3, ha='center', transform=ax_t.transAxes)
        ax_t.text(0.5, 0.42, '30 визуализаций на основе реальных данных проекта Fundament RF', fontsize=12,
                 color=TEXT_DIM, ha='center', transform=ax_t.transAxes)
        pdf.savefig(fig_t, facecolor=BG_COLOR)
        plt.close(fig_t)

        for name, fig in slides:
            pdf.savefig(fig, facecolor=BG_COLOR)
            plt.close(fig)

    print(f"✅ PDF: {pdf_path}")

    # ─── Build PPTX ────────────────────────────────────
    print(f"\nСборка PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    from pptx.util import Pt as PtSize
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Учебник институционального трейдинга"
    p.font.size = PtSize(36)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "От Price Action до рыночной микроструктуры"
    p2.font.size = PtSize(18)
    p2.font.color.rgb = RGBColor(0xff, 0xea, 0x00)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "30 визуализаций на основе реальных данных проекта Fundament RF"
    p3.font.size = PtSize(12)
    p3.font.color.rgb = RGBColor(0xcc, 0xcc, 0xdd)
    p3.alignment = PP_ALIGN.CENTER

    for name, fig in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

        # Save figure to PNG in memory
        img_buf = io.BytesIO()
        fig.savefig(img_buf, format='png', dpi=120, bbox_inches='tight',
                   facecolor=BG_COLOR, edgecolor='none')
        img_buf.seek(0)

        slide.shapes.add_picture(img_buf, Inches(0), Inches(0), Inches(16), Inches(9))

    pptx_path = OUTPUT_DIR / 'Презентация_учебника.pptx'
    prs.save(str(pptx_path))
    print(f"✅ PPTX: {pptx_path} ({len(slides)+1} слайдов)")

    # Create index file
    index_path = OUTPUT_DIR / 'Презентация_учебника.md'
    with open(index_path, 'w') as f:
        f.write("# Презентация учебника институционального трейдинга\n\n")
        f.write("## Содержание слайдов\n\n")
        for i, (name, _) in enumerate(slides, 1):
            f.write(f"{i:2d}. {name}\n")
    print(f"✅ Index: {index_path}")

    print(f"\n{'='*60}")
    print(f"ГОТОВО!")
    print(f"  PPTX: {pptx_path}")
    print(f"  PDF:  {pdf_path}")
    print(f"{'='*60}")

if __name__ == '__main__':
    main()
